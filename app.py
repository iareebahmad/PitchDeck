import streamlit as st
import openai
from dotenv import load_dotenv
import os
import fitz  # PyMuPDF
from pptx import Presentation

# Load API Key from .env
load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")

st.set_page_config(page_title='Pitch Deck Analysis')
st.title("Pitch Deck Analysis")


uploaded_file = st.file_uploader("Upload your Pitch Deck (PDF or PPTX)", type=["pdf", "pptx"])

def extract_text_from_pdf(file):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

if uploaded_file:
    st.success("✅ File uploaded!")

    if uploaded_file.type == "application/pdf":
        pitch_text = extract_text_from_pdf(uploaded_file)
    else:
        pitch_text = extract_text_from_pptx(uploaded_file)

    if st.button("Analyze Pitch Deck"):
        prompt = f"""
        You are a venture capital expert.  
        Based on the following pitch deck text, analyze the 10 aspects:  
        Solution, Funding Ask, Market Size, Business Model, Team Credentials, Traction Metrics, Problem Statement, Competitive Landscape, Financial Projections, Go To Market Strategy.

        Finally, provide an Investment Readiness Score in percentage (0-100%) along with a brief reasoning in not more than 80 words.

        Pitch Deck Text:
        {pitch_text}

        Output format:
        Investment Readiness: X%
        Reason: [Brief explanation]
        """

        with st.spinner("Analyzing..."):
            try:
                response = openai.ChatCompletion.create(
                    model="gpt-4o-mini",
                    messages=[
                        {"role": "system", "content": "You are a venture capital expert."},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.7,
                    max_tokens=200,
                )

                result = response['choices'][0]['message']['content']
                st.success("Analysis Complete!")
                st.markdown(f"### 📈 Result:\n{result}")

            except Exception as e:
                st.error(f"❌ Error: {e}")
